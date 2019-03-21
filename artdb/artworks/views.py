import logging
import operator
import os
from datetime import datetime
from functools import reduce
from io import BytesIO

from dal import autocomplete
from django.conf import settings
from django.contrib.auth.decorators import login_required, permission_required
from django.contrib.auth.models import User
from django.core.paginator import Paginator, EmptyPage, PageNotAnInteger
from django.db.models import Q, ExpressionWrapper, BooleanField, Case, Value, IntegerField, When
from django.http import JsonResponse, HttpResponse, HttpResponseForbidden
from django.shortcuts import _get_queryset
from django.shortcuts import render, get_object_or_404, redirect
from django.template.defaultfilters import slugify
from django.utils.decorators import method_decorator
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt
from rest_framework.exceptions import APIException

from .forms import ArtworkForm, ArtworkCollectionForm
from .models import Artwork, Location, Keyword, Artist, ArtworkCollection, ArtworkCollectionMembership
from .serializers import ArtworkSerializer, CollectionSerializer

logger = logging.getLogger(__name__)


def get_object_or_APIException(klass, *args, **kwargs):
    """
    Similar to get_object_or_404, but raises APIException.
    """
    queryset = _get_queryset(klass)
    if not hasattr(queryset, 'get'):
        klass__name = klass.__name__ if isinstance(klass, type) else klass.__class__.__name__
        raise ValueError(
            "First argument to get_object_or_APIException() must be a Model, Manager, "
            "or QuerySet, not '%s'." % klass__name
        )
    try:
        return queryset.get(*args, **kwargs)
    except queryset.model.DoesNotExist:
        klass__name = klass.__name__ if isinstance(klass, type) else klass.__class__.__name__
        message = 'Could not find %s %s' % (klass__name, kwargs)
        logger.error(message)
        raise APIException(message)


@login_required
def artworks_list(request):
    """
    Render the thumbnailbrowser.
    """
    query_search = request.GET.get('search')
    query_search_type = request.GET.get('searchtype')
    query_artwork_title = request.GET.get('title')
    query_artist_name = request.GET.get('artist')
    query_keyword = request.GET.get('keyword')
    query_date_from = request.GET.get('date_from')
    query_date_to = request.GET.get('date_to')
    query_location_of_creation = request.GET.get('location_of_creation')
    query_location_current = request.GET.get('location_current')
    q_objects = Q()
    context = {}
    expert_search = False

    def get_expert_queryset_list():
        expert_list = Artwork.objects.filter(published=True)
        if query_location_of_creation:
            locations = Location.objects.filter(name__istartswith=query_location_of_creation)
            locations_plus_descendants = Location.objects.get_queryset_descendants(locations, include_self=True)
            q_objects.add(Q(location_of_creation__in=locations_plus_descendants), Q.AND)
        if query_location_current:
            locations = Location.objects.filter(name__istartswith=query_location_current)
            locations_plus_descendants = Location.objects.get_queryset_descendants(locations, include_self=True)
            q_objects.add(Q(location_current__in=locations_plus_descendants), Q.AND)
        if query_artist_name:
            terms = [term.strip() for term in query_artist_name.split()]
            for term in terms:
                q_objects.add((Q(artists__name__icontains=term) | Q(artists__synonyms__icontains=term)), Q.AND)
        if query_keyword:
            keywords = Keyword.objects.filter(name__icontains=query_keyword)
            q_objects.add(Q(keywords__in=keywords), Q.AND)
        if query_date_from:
            try:
                year = int(query_date_from)
                q_objects.add(Q(date_year_from__gte=year), Q.AND)
            except ValueError as err:
                logger.error(err)
                return []
        if query_date_to:
            try:
                year = int(query_date_to)
                q_objects.add(Q(date_year_to__lte=year), Q.AND)
            except ValueError as err:
                logger.error(err)
                return []
        if query_artwork_title:
            title_contains = (Q(title__icontains=query_artwork_title) |
                              Q(title_english__icontains=query_artwork_title))
            title_starts_with = (Q(title__istartswith=query_artwork_title) |
                                 Q(title_english__istartswith=query_artwork_title))
            # order results by startswith match. see: https://stackoverflow.com/a/48409962
            expert_list = expert_list.filter(title_contains)
            is_match = ExpressionWrapper(title_starts_with, output_field=BooleanField())
            expert_list = expert_list.annotate(starts_with_title=is_match)
            expert_list = expert_list.filter(q_objects).order_by('-starts_with_title', 'location_of_creation')
        else:
            expert_list = expert_list.filter(q_objects).order_by('title', 'location_of_creation')
        return expert_list.distinct()

    def get_basic_queryset_list():
        if query_search:

            def get_artists(term):
                return Artist.objects.filter(Q(name__istartswith=term) | Q(name__icontains=' ' + term))

            def get_keywords(term):
                return Keyword.objects.filter(Q(name__istartswith=term) | Q(name__istartswith=' ' + term))

            terms = [term.strip() for term in query_search.split()]
            basic_list = (Artwork.objects.annotate(
                rank=Case(
                    When(Q(title__iexact=query_search), then=Value(1)),
                    When(Q(title_english__iexact=query_search), then=Value(1)),
                    When(Q(artists__in=get_artists(query_search)), then=Value(2)),
                    When(Q(title__istartswith=query_search), then=Value(3)),
                    When(Q(title_english__istartswith=query_search), then=Value(3)),
                    When(reduce(operator.or_, (Q(artists__in=get_artists(term)) for term in terms)), then=Value(4)),
                    When(reduce(operator.or_, (Q(title__istartswith=term) for term in terms)), then=Value(5)),
                    When(reduce(operator.or_, (Q(title_english__istartswith=term) for term in terms)), then=Value(5)),
                    When(reduce(operator.or_, (Q(title__icontains=' ' + term) for term in terms)), then=Value(6)),
                    When(reduce(operator.or_, (Q(title_english__icontains=' ' + term) for term in terms)), then=Value(6)),
                    When(reduce(operator.or_, (Q(title__icontains=term) for term in terms)), then=Value(7)),
                    When(reduce(operator.or_, (Q(title_english__icontains=term) for term in terms)), then=Value(7)),
                    When(reduce(operator.or_, (Q(location_of_creation__name__istartswith=term) for term in terms)), then=Value(10)),
                    When(reduce(operator.or_, (Q(keywords__in=get_keywords(term)) for term in terms)), then=Value(11)),
                    default=Value(99),
                    output_field=IntegerField(),
                ))
                .filter(published=True)
                .exclude(rank=99)
                .distinct('id', 'rank', 'title',)
                .order_by('rank', 'title',))
        else:
            # what the user gets, when she isn't using the search at all
            basic_list = Artwork.objects.filter(published=True).order_by('-updated_at', 'title')
        return basic_list

    if query_search_type == 'expert':
        queryset_list = get_expert_queryset_list()
        expert_search = True
    else:
        queryset_list = get_basic_queryset_list()

    paginator = Paginator(queryset_list, 40)  # show 40 artworks per page
    page_nr = request.GET.get('page')
    try:
        artworks = paginator.get_page(page_nr)
    except PageNotAnInteger:
        artworks = paginator.page(1)
    except EmptyPage:
        artworks = paginator.page(paginator.num_pages)

    context['artworks'] = artworks
    context['query_search'] = query_search
    context['query_title'] = query_artwork_title
    context['query_artist'] = query_artist_name
    context['query_keyword'] = query_keyword
    context['query_date_from'] = query_date_from
    context['query_date_to'] = query_date_to
    context['query_location_of_creation'] = query_location_of_creation
    context['query_location_current'] = query_location_current
    context['expert_search'] = expert_search
    return render(request, 'artwork/thumbnailbrowser.html', context)


@login_required
def details(request, id=None):
    """
    Return artwork details in json format.
    """
    try:
        artwork = get_object_or_APIException(Artwork, id=id)
    except APIException:
        return JsonResponse(status=404, data={'status': 'false', 'message': 'Could not get artwork details'})
    serializer = ArtworkSerializer(artwork)
    return JsonResponse(serializer.data)


@login_required
def artwork_detail_overlay(request, id=None):
    """
    Render an overlay showing a large version of the image and the artwork's details.
    """
    artwork = get_object_or_404(Artwork, id=id)
    context = {
        'artwork': artwork,
        'is_staff': request.user.is_staff,
    }
    return render(request, 'artwork/artwork_detail_overlay.html', context)


@permission_required('artworks.change_artwork')
def artwork_edit(request, id):
    """
    Render an overlay showing the editable fields of an artwork.
    """
    artwork = get_object_or_404(Artwork, id=id)
    if request.method == "POST":
        form = ArtworkForm(request.POST, request.FILES, instance=artwork)
        if form.is_valid():
            updated_artwork = form.save(commit=False)
            updated_artwork.updated_at = datetime.now()
            updated_artwork.save()
            return HttpResponse("<script>window.location=document.referrer;</script>")
    context = {
        'form': ArtworkForm(instance=artwork),
        'id': artwork.id,
        'image_original': artwork.image_original,
    }
    return render(request, 'artwork/artwork_edit_overlay.html', context)


@login_required
def artwork_collect(request, id):
    """
    Add or remove an artwork from/to a collection.
    """
    if request.method == 'GET':
        artwork = get_object_or_404(Artwork, id=id)
        context = {}
        qs = ArtworkCollection.objects.all()
        collections = qs.filter(user=request.user).order_by('-created_at')
        context['collections'] = collections
        context['artwork'] = artwork
        return render(request, 'artwork/artwork_collect_overlay.html', context)
    if request.method == 'POST':
        try:
            artwork = get_object_or_APIException(Artwork, id=request.POST['artwork-id'])
        except APIException:
            return JsonResponse(status=404, data={'status': 'false', 'message': 'Could not get artwork'})
        if request.POST['action'] == 'addCollection':
            col_title = request.POST['collection-title']
            if col_title:
                try:
                    u = get_object_or_APIException(User, id=request.user.id)
                    newcol = ArtworkCollection.objects.create(title=col_title, user=u)
                    ArtworkCollectionMembership.objects.create(collection=newcol, artwork=artwork)
                    return JsonResponse({'action': 'reload'})
                except APIException:
                    return JsonResponse(status=404, data={'status': 'false', 'message': 'Could not get user'})
            else:
                return JsonResponse({'error': 'collection title missing'}, status=500)
        else:
            try:
                col = get_object_or_APIException(ArtworkCollection, id=request.POST['collection-id'])
            except APIException:
                return JsonResponse(status=404, data={'status': 'false', 'message': 'Could not get collection'})
            # users can only manipulate their own collections via this view
            if request.user == col.user:
                if request.POST['action'] == 'add':
                    try:
                        ArtworkCollectionMembership.objects.create(collection=col, artwork=artwork)
                    except APIException:
                        return JsonResponse(status=500, data={'status': 'false', 'message': 'Could not add artwork to collection'})
                    return JsonResponse({'action': 'added'})
                if request.POST['action'] == 'remove':
                    try:
                        artworkColMem = get_object_or_APIException(ArtworkCollectionMembership, artwork=artwork)
                        artworkColMem.remove()
                        return JsonResponse({'action': 'removed'})
                    except APIException:
                        return JsonResponse(
                            data={'status': 'false', 'message': 'Could not remove artwork from collection'},
                            status=500,
                        )
        return JsonResponse(status=500, data={'status': 'false', 'message': 'Could not manipulate collection'})


@login_required
def collection(request, id=None):
    """
    GET: Render all artwork thumbnails of a single collection.
    POST: move artworks within collection; connect or disconnect them
    """
    if request.method == 'GET':
        col = get_object_or_404(ArtworkCollection, id=id)
        context = {
            'title': col.title,
            'id': col.id,
            'created_by_username': col.user.get_username(),
            'created_by_fullname': col.user.get_full_name(),
            'created_by_userid': col.user.id,
            'memberships': col.artworkcollectionmembership_set.all(),
            'collections': ArtworkCollection.objects.filter(user__groups__name='editor').exclude(user=request.user),
            'my_collections': ArtworkCollection.objects.filter(user=request.user),
        }
        return render(request, 'artwork/collection.html', context)
    if request.method == 'POST':
        # users can only manipulate their own collections via this view
        try:
            col = get_object_or_APIException(ArtworkCollection, id=id)
        except APIException:
            return JsonResponse(status=404, data={'status': 'false', 'message': 'Could not find artwork collection'})
        if request.user.id != col.user.id:
            return JsonResponse(status=403, data={'status': 'false', 'message': 'Permission needed'})
        if request.POST['action'] == 'left' or request.POST['action'] == 'right':
            if 'membership-id' in request.POST:
                try:
                    membership = get_object_or_APIException(ArtworkCollectionMembership, id=request.POST['membership-id'], collection=col)
                    # move artwork left
                    if request.POST['action'] == 'left':
                        membership.move_left()
                        return JsonResponse({'message': 'moved left'})
                    # move artwork right
                    if request.POST['action'] == 'right':
                        membership.move_right()
                        return JsonResponse({'message': 'moved right'})
                except APIException:
                    logger.error("Could not move artwork membership: %s", request.POST['membership-id'])
            return JsonResponse(status=500, data={'status': 'false', 'message': 'Could not move artwork membership'})
        elif request.POST['action'] == 'connect' or request.POST['action'] == 'disconnect':
            # user wants to connect or disconnect the artwork
            try:
                left_member = get_object_or_APIException(
                    ArtworkCollectionMembership,
                    id=request.POST['member-left'],
                    collection=col,
                )
                right_member = get_object_or_APIException(
                    ArtworkCollectionMembership,
                    id=request.POST['member-right'],
                    collection=col,
                )
                if request.POST['action'] == 'connect':
                    if left_member.connect(right_member):
                        return JsonResponse({'message': 'connected'})
                if request.POST['action'] == 'disconnect':
                    if left_member.disconnect(right_member):
                        return JsonResponse({'message': 'disconnected'})
            except APIException:
                logger.error(
                    "Could not do action '%s' on artwork membership '%s'",
                    request.POST['action'],
                    request.POST['membership-id'],
                )
            return JsonResponse(
                data={'status': 'false', 'message': 'Could not connect/disconnect artwork membership'},
                status=500,
            )
        return JsonResponse(data={'status': 'false', 'message': 'Invalid post action'}, status=404)


@login_required
def collection_edit(request, id):
    """
    Render an overlay showing the editable fields of a collection.
    """
    collection = get_object_or_404(ArtworkCollection, id=id)
    if request.user.id is not collection.user.id:
        # users can only manipulate their own collections via this view
        return HttpResponseForbidden()
    if request.method == "POST":
        form = ArtworkCollectionForm(request.POST, instance=collection)
        if form.is_valid():
            form.save()
            return redirect('collection', id=id)
    context = {
        'form': ArtworkCollectionForm(instance=collection),
        'collection': collection,
    }
    return render(request, 'artwork/collection_edit_overlay.html', context)


@login_required
def collection_delete(request, id):
    """
    Delete a collection.
    """
    if request.method == "POST":
        try:
            collection = get_object_or_APIException(ArtworkCollection, id=id)
            if request.user.id is collection.user.id:
                # users can only manipulate their own collections via this view
                collection.delete()
                return redirect('collections-list')
            else:
                return JsonResponse(status=403, data={'status': 'false', 'message': 'Permission needed'})
        except APIException:
            return JsonResponse(status=500, data={'status': 'false', 'message': 'Could not delete collection'})
    else:
        return redirect('collection', id=id)


@login_required
def collection_json(request, id=None):
    """
    Return collection data in json format.
    """
    try:
        col = get_object_or_APIException(ArtworkCollection, id=id)
    except APIException:
        return JsonResponse(status=404, data={'status': 'false', 'message': 'Could not get collection'})
    serializer = CollectionSerializer(col)
    return JsonResponse(serializer.data)


@login_required
def collections_list(request):
    """
    Render a list of all collections.
    """
    context = {
        'collections': ArtworkCollection.objects.filter(user__groups__name='editor').exclude(user=request.user),
        'my_collections': ArtworkCollection.objects.filter(user=request.user),
    }
    return render(request, 'artwork/collections_list.html', context)


@method_decorator(login_required, name='dispatch')
class ArtistAutocomplete(autocomplete.Select2QuerySetView):
    """
    Return dal suggestions for the artist input field.
    """
    def get_queryset(self):
        qs = Artist.objects.all().order_by('name')
        if self.q:
            return qs.filter(name__istartswith=self.q)
        else:
            return Artist.objects.none()


@method_decorator(login_required, name='dispatch')
class ArtworkAutocomplete(autocomplete.Select2QuerySetView):
    """
    Return dal suggestions for the basic search.
    Suggest the first 4 artworks with matching titles.
    """
    def get_queryset(self):
        qs = Artwork.objects.all()
        if self.q:
            qs = qs.filter(title__icontains=self.q).order_by('title')
            # order results by startswith match. see: https://stackoverflow.com/a/48409962
            expression = Q(title__istartswith=self.q) | Q(title__icontains=' ' + self.q)
            is_match = ExpressionWrapper(expression, output_field=BooleanField())
            qs = qs.annotate(title_starts_with=is_match)
            qs = qs.distinct('title', 'title_starts_with')
            qs = qs.order_by('-title_starts_with', 'title')[:4]
            return qs
        else:
            return Artwork.objects.none()
            
    def get_result_value(self, result):
        """Return the value of a result."""
        return result.title


@method_decorator(login_required, name='dispatch')
class KeywordAutocomplete(autocomplete.Select2QuerySetView):
    """
    Return dal suggestions for the artwork's keywords input field.
    """
    def get_queryset(self):
        qs = Keyword.objects.all().order_by('name')
        if self.q:
            return qs.filter(name__icontains=self.q)
        else:
            return Keyword.objects.none()


def collection_download_as_pptx_en(request, id=None):
    return collection_download_as_pptx(request, id, 'en')


def collection_download_as_pptx_de(request, id=None):
    return collection_download_as_pptx(request, id, 'de')


def collection_download_as_pptx(request, id=None, language='de'):
    """
    Return a downloadable powerpoint presentation of the collection
    """
    def get_new_slide():
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 30, 30)
        return slide

    def add_description(slide, description, width, left):
        shapes = slide.shapes
        top = prs.slide_height - textbox_height - padding
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, textbox_height)
        shape.fill.background()
        shape.line.fill.background()
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = description
        font = run.font
        font.size = Pt(30)

    def add_slide_with_one_picture(artwork, padding):
        img_relative_path = artwork.image_original.thumbnail['1881x933'].name
        img_path = os.path.join(settings.MEDIA_ROOT, img_relative_path)
        slide = get_new_slide()
        add_picture_to_slide(slide, img_path, padding, 'center')
        picture_width = prs.slide_width - (padding * 2)
        add_description(slide, artwork.get_short_description(language), picture_width, padding)

    def add_slide_with_two_pictures(artwork_left, artwork_right, padding):
        img_relative_path_left = artwork_left.image_original.thumbnail['576x933'].name
        img_path_left = os.path.join(settings.MEDIA_ROOT, img_relative_path_left)
        img_relative_path_right = artwork_right.image_original.thumbnail['576x933'].name
        img_path_right = os.path.join(settings.MEDIA_ROOT, img_relative_path_right)
        slide = get_new_slide()
        add_picture_to_slide(slide, img_path_left, padding, 'left')
        add_picture_to_slide(slide, img_path_right, padding, 'right')
        text_width = int((prs.slide_width - (padding * 2) - distance_between)/2)
        add_description(slide, artwork_left.get_short_description(language), text_width, padding)
        left = padding + text_width + distance_between
        add_description(slide, artwork_right.get_short_description(language), text_width, left)

    def add_picture_to_slide(slide, img_path, padding, position):
        pic = slide.shapes.add_picture(img_path, 0, padding)
        image_width = pic.image.size[0]
        image_height = pic.image.size[1]
        aspect_ratio = image_width / image_height
        
        # calculate width and height
        if position == 'center':
            picture_max_width = int(prs.slide_width - (padding * 2))
            space_aspect_ratio = picture_max_width / picture_max_height
            if aspect_ratio < space_aspect_ratio:
                pic.height = picture_max_height
                pic.width = int(picture_max_height * aspect_ratio)
            else:
                pic.width = picture_max_width
                pic.height = int(picture_max_width / aspect_ratio)
                pic.top = padding + int((picture_max_height - pic.height) / 2)
        else:
            picture_max_width = int((prs.slide_width - (padding * 2) - distance_between)/2)
            space_aspect_ratio = picture_max_width / picture_max_height
            if aspect_ratio < space_aspect_ratio:
                pic.height = picture_max_height
                pic.width = int(picture_max_height * aspect_ratio)
            else:
                pic.width = picture_max_width
                pic.height = int(picture_max_width / aspect_ratio)
                pic.top = padding + int((picture_max_height - pic.height) / 2)

        # position the image left/right
        if position == 'center':
            pic.left = int((prs.slide_width - pic.width) / 2)
        if position == 'left':
            if image_height < image_width:
                pic.left = int(padding)
            else:
                pic.left = padding + int((picture_max_width - pic.width)/2)
        if position == 'right':
            if image_height < image_width:
                pic.left = padding + picture_max_width + distance_between
            else:
                pic.left = padding + picture_max_width + distance_between + int((picture_max_width - pic.width)/2)

    # define the presentation dimensions
    prs = Presentation()
    prs.slide_width = 24384000   # taken from Keynote 16:9 pptx
    prs.slide_height = 13716000  # taken from Keynote 16:9 pptx
    padding = int(prs.slide_width / 100)
    textbox_height = prs.slide_height / 10
    picture_max_height = int(prs.slide_height - (padding * 2) - textbox_height)
    distance_between = padding * 2

    try:
        col = get_object_or_APIException(ArtworkCollection, id=id)
    except APIException:
        logger.error('Could not create powerpoint file. Collection missing.')
        return

    memberships = col.artworkcollectionmembership_set.all()
    connected_member = None
    for membership in memberships:
        if membership.connected_with:
            if not connected_member:
                connected_member = membership
            else:
                add_slide_with_two_pictures(connected_member.artwork, membership.artwork, padding)
                connected_member = None
        else:
            add_slide_with_one_picture(membership.artwork, padding)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    response = HttpResponse(
        output.read(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )
    response['Content-Disposition'] = 'attachment; filename="' + slugify(col.title) + '.pptx"'
    output.close()

    return response
