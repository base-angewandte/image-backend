import logging
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt

from django.conf import settings
from django.http import HttpResponse
from django.template.defaultfilters import slugify
from django.utils.translation import gettext_lazy as _

from .models import Album, Artwork

logger = logging.getLogger(__name__)


class ExportError(Exception):
    pass


def album_download_as_pptx(album_id, language='en', return_raw=False):
    """Return a downloadable PowerPoint presentation of the album."""

    try:
        album = Album.objects.get(id=album_id)
    except Album.DoesNotExist as dne:
        logger.warning('Could not create powerpoint file. Album missing.')
        raise ExportError(_('Album does not exist')) from dne

    # define the presentation dimensions
    prs = Presentation()
    prs.slide_width = 24384000  # taken from Keynote 16:9 pptx
    prs.slide_height = 13716000  # taken from Keynote 16:9 pptx

    prs_padding = int(prs.slide_width / 96)  # full HD: 1920px/96 = 20px
    textbox_height = prs.slide_height / 10
    picture_max_height = int(prs.slide_height - (prs_padding * 2) - textbox_height)
    distance_between = prs_padding * 2

    def add_run_to_paragraph(paragraph, text, style=None):
        run = paragraph.add_run()
        run.text = text
        font = run.font
        font.size = Pt(36)
        font.color.theme_color = MSO_THEME_COLOR.TEXT_1
        if style == 'strikethrough':
            font._element.attrib['strike'] = 'sngStrike'
            font._element.attrib['baseline'] = '-25000'

    def get_new_slide():
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(217, 217, 217)
        return slide

    def add_description(slide, artwork, width, left):
        shapes = slide.shapes
        top = prs.slide_height - textbox_height - prs_padding
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, textbox_height)
        shape.fill.background()
        shape.line.fill.background()
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]

        # apply discriminatory terms styling
        discriminatory_terms = artwork.get_discriminatory_terms_list(
            order_by_length=True,
        )
        description = artwork.get_short_description(language)

        # limit discriminatory terms to the ones in description
        discriminatory_terms = [
            term
            for term in discriminatory_terms
            if description.lower().find(term.lower()) != -1
        ]

        # Process the text by finding all occurrences of the terms

        # Track the position within the description
        index = 0

        # while a regexp based approach could be desirable, we still need to walk
        # through the whole description index-wise, because we cannot simply replace
        # but need to build the whole paragraph with runs
        while index < len(description):
            found_term = None
            found_position = len(description)

            # find the lowest position of any matched term
            for term in discriminatory_terms:
                pos = description.lower().find(term.lower(), index)
                if pos != -1 and pos < found_position:
                    found_term = term
                    found_position = pos

            if not found_term:
                add_run_to_paragraph(p, description[index:])
                break

            # Add the text before the found term
            if index < found_position:
                add_run_to_paragraph(p, description[index:found_position])

            # Add the first letter of the found term in normal style
            add_run_to_paragraph(p, description[found_position])

            # Apply style 'strikethrough' to rest of the term
            add_run_to_paragraph(
                p,
                description[found_position : found_position + len(found_term)],
                style='strikethrough',
            )

            # Move the index forward after processing the found term
            index = found_position + len(found_term)

    def add_picture_to_slide(slide, img_path: Path, padding, position):
        pic = slide.shapes.add_picture(img_path.as_posix(), 0, padding)
        image_width = pic.image.size[0]
        image_height = pic.image.size[1]
        aspect_ratio = image_width / image_height

        # calculate width and height
        if position == 'center':
            picture_max_width = int(prs.slide_width - (padding * 2))
            space_aspect_ratio = picture_max_width / picture_max_height

            if aspect_ratio < space_aspect_ratio:
                pic.height = picture_max_height
                pic.width = int(picture_max_height * aspect_ratio)
            else:
                pic.width = picture_max_width
                pic.height = int(picture_max_width / aspect_ratio)
                pic.top = padding + int((picture_max_height - pic.height) / 2)

        else:
            picture_max_width = int(
                (prs.slide_width - (padding * 2) - distance_between) / 2,
            )
            space_aspect_ratio = picture_max_width / picture_max_height

            if aspect_ratio < space_aspect_ratio:
                pic.height = picture_max_height
                pic.width = int(picture_max_height * aspect_ratio)
            else:
                pic.width = picture_max_width
                pic.height = int(picture_max_width / aspect_ratio)
                pic.top = padding + int((picture_max_height - pic.height) / 2)

        # position the image center/left/right
        match position:
            case 'center':
                pic.left = int((prs.slide_width - pic.width) / 2)
            case 'left':
                if image_height < image_width:
                    pic.left = int(padding)
                else:
                    pic.left = padding + int((picture_max_width - pic.width) / 2)
            case 'right':
                if image_height < image_width:
                    pic.left = padding + picture_max_width + distance_between
                else:
                    pic.left = (
                        padding
                        + picture_max_width
                        + distance_between
                        + int((picture_max_width - pic.width) / 2)
                    )

    def add_slide(artworks: list, padding):
        # slide with one image
        if len(artworks) == 1:
            thumbnail_size = '1880x933'  # 1920-20-20 = 1880

            artwork = artworks[0]

            img_relative_path = artwork.image_fullsize.thumbnail[thumbnail_size].name
            img_path = settings.MEDIA_ROOT_PATH / img_relative_path

            slide = get_new_slide()

            add_picture_to_slide(slide, img_path, padding, 'center')

            text_width = prs.slide_width - (padding * 2)

            add_description(slide, artwork, text_width, padding)

        # slide with two images
        elif len(artworks) == 2:
            thumbnail_size = '920x933'  # (1920/2)-20-20 = 920

            artwork_left = artworks[0]
            artwork_right = artworks[1]

            img_relative_path_left = artwork_left.image_fullsize.thumbnail[
                thumbnail_size
            ].name
            img_path_left = settings.MEDIA_ROOT_PATH / img_relative_path_left

            img_relative_path_right = artwork_right.image_fullsize.thumbnail[
                thumbnail_size
            ].name
            img_path_right = settings.MEDIA_ROOT_PATH / img_relative_path_right

            slide = get_new_slide()

            add_picture_to_slide(slide, img_path_left, padding, 'left')
            add_picture_to_slide(slide, img_path_right, padding, 'right')

            text_width = int((prs.slide_width - (padding * 2) - distance_between) / 2)

            add_description(slide, artwork_left, text_width, padding)

            left = padding + text_width + distance_between

            add_description(slide, artwork_right, text_width, left)

    slides = album.slides

    # TODO: for now we just drop artworks which do not exist any more from the slides
    #   in a future feature we need to discuss whether there should be some information left, that there was
    #   an artwork but got deleted, and whether we should retain some artwork title in that case, or just
    #   display a blank). technically, we could add an Album.repair_slides() method which handles this

    if slides:
        for slide in slides:
            try:
                if len(slide['items']) <= 2:
                    artworks = []
                    for item in slide['items']:
                        try:
                            artwork = Artwork.objects.get(
                                id=item.get('id'),
                                published=True,
                            )
                            artworks.append(artwork)
                        except Artwork.DoesNotExist:
                            pass

                    if artworks:
                        add_slide(artworks, prs_padding)

                else:
                    raise ExportError(
                        _('Album contains slides with more than 2 artworks'),
                    )
            except FileNotFoundError as fnfe:
                raise ExportError(
                    _('At least one image file can not be found'),
                ) from fnfe

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    if return_raw:
        return output

    response = HttpResponse(
        output.read(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={
            'Content-Disposition': f'attachment; filename={slugify(album.title)}.pptx',
        },
    )

    output.close()

    return response
